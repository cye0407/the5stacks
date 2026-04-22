"""
Build Phase 1 Video 2: The History + The Lenses — PowerPoint deck (v2)
Style-matched to 5_stacks_physics_video_deck-v-3.1
- Pure black backgrounds
- Large bold uppercase white headers
- Minimal text per slide
- Red/green accents
- Lots of breathing room
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors matched to v3.1
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA0, 0xA0, 0xA0)
DARK_GRAY = RGBColor(0x70, 0x70, 0x70)
GREEN = RGBColor(0x22, 0x7C, 0x5C)
RED = RGBColor(0xDC, 0x26, 0x26)
AMBER = RGBColor(0xD9, 0x77, 0x06)
BLUE_BOX = RGBColor(0x4A, 0x72, 0x9A)  # from stabilizing cycle
RED_BOX = RGBColor(0xC0, 0x50, 0x50)   # from amplifying cycle
YELLOW = RGBColor(0xFB, 0xBF, 0x24)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
FONT = "Segoe UI"


def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK


def text(slide, left, top, width, height, content, size=28,
         color=WHITE, bold=False, align=PP_ALIGN.CENTER, italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = FONT
    p.alignment = align
    return txBox


def multiline(slide, left, top, width, height, lines, size=24, color=WHITE,
              align=PP_ALIGN.CENTER, spacing=12, bolds=None, italics=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(spacing)
        p.alignment = align
        if bolds and i < len(bolds):
            p.font.bold = bolds[i]
        if italics and i < len(italics):
            p.font.italic = italics[i]
    return txBox


def shape(slide, stype, left, top, width, height, fill=None, line=None, line_width=2):
    s = slide.shapes.add_shape(stype, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(line_width)
    else:
        s.line.fill.background()
    return s


# ============================================================
# SLIDE 1: Title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 1.5, 14, 2,
     "THE WORD SUSTAINABILITY",
     size=54, color=WHITE, bold=True)

text(s, 1, 3.2, 14, 2,
     "IS GERMAN",
     size=54, color=WHITE, bold=True)

text(s, 1, 5.5, 14, 1,
     "It's 300 years old.",
     size=28, color=GRAY)

text(s, 1, 6.5, 14, 1,
     "It had nothing to do with environmentalism.",
     size=28, color=GRAY)

text(s, 1, 8, 14, 0.8,
     "THE 5 STACKS  |  Video 2 of 5",
     size=16, color=DARK_GRAY)


# ============================================================
# SLIDE 2: Before the word — Indigenous
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 1, 14, 2,
     "THOUSANDS OF YEARS",
     size=54, color=WHITE, bold=True)

text(s, 1, 2.7, 14, 1.5,
     "BEFORE THE WORD EXISTED",
     size=54, color=WHITE, bold=True)

multiline(s, 2, 5, 12, 3, [
    "Controlled burns. Rotational farming. Seasonal harvesting.",
    "",
    "Every continent. Every culture.",
    "Not called 'sustainability.' Called survival.",
], size=24, color=GRAY, spacing=10)


# ============================================================
# SLIDE 3: Manage or starve
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 2, 14, 2,
     "YOU MANAGE RESOURCES",
     size=48, color=WHITE, bold=True)

text(s, 1, 4, 14, 1.5,
     "OR YOU DON'T EAT NEXT YEAR",
     size=48, color=WHITE, bold=True)

multiline(s, 2, 6.5, 12, 2, [
    "No frameworks. No reports. No consultants.",
    "Just intergenerational common sense.",
], size=24, color=GRAY, spacing=10)


# ============================================================
# SLIDE 4: Edo Japan
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 0.8, 14, 1.5,
     "EDO JAPAN",
     size=54, color=WHITE, bold=True)

text(s, 1, 2.3, 14, 1,
     "1603 \u2013 1868",
     size=28, color=GRAY)

text(s, 1, 3.8, 14, 1,
     "250 years of near-total isolation.",
     size=28, color=GRAY)

text(s, 1, 4.8, 14, 1,
     "Nothing wasted because nothing could be replaced.",
     size=28, color=GRAY)

# Simple flow: clothing > rags > paper
items = ["New clothing", "Worn clothing", "Rags", "Paper"]
x_start = 1.5
box_w = 2.8
gap = 0.5
y = 6.5

for i, item in enumerate(items):
    x = x_start + i * (box_w + gap + 0.3)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, 1.2, line=GRAY, line_width=1)
    text(s, x, y + 0.2, box_w, 0.8, item, size=20, color=WHITE)
    if i < len(items) - 1:
        shape(s, MSO_SHAPE.RIGHT_ARROW, x + box_w + 0.05, y + 0.3, 0.4, 0.5, fill=GRAY)

text(s, 1, 8.2, 14, 0.8,
     "Circularity by necessity \u2014 not by certification.",
     size=20, color=DARK_GRAY, italic=True)


# ============================================================
# SLIDE 5: Medieval Commons
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 1, 14, 2,
     "MEDIEVAL EUROPEAN",
     size=54, color=WHITE, bold=True)

text(s, 1, 2.7, 14, 1.5,
     "COMMONS",
     size=54, color=WHITE, bold=True)

text(s, 1, 4.5, 14, 1,
     "Communities set limits on extraction.",
     size=28, color=GRAY)

text(s, 1, 5.5, 14, 1,
     "Not because they were green. Because overharvesting destroyed everyone.",
     size=24, color=GRAY)

# Three pillars
pillars = [
    ("Shared grazing", "Limits per household"),
    ("Forestry rules", "Cut one, plant two"),
    ("Water rights", "Don't poison downstream"),
]

x_start = 2
box_w = 3.5
for i, (title, desc) in enumerate(pillars):
    x = x_start + i * (box_w + 0.5)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, 7, box_w, 1.5, line=GREEN, line_width=1)
    text(s, x, 7.1, box_w, 0.7, title, size=20, color=GREEN, bold=True)
    text(s, x, 7.7, box_w, 0.6, desc, size=16, color=GRAY)


# ============================================================
# SLIDE 6: Carlowitz 1713
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 0.8, 14, 2,
     "1713",
     size=96, color=WHITE, bold=True)

text(s, 1, 3, 14, 1,
     "HANS CARL VON CARLOWITZ",
     size=36, color=WHITE, bold=True)

text(s, 1, 4.2, 14, 1,
     "Saxon mining administrator",
     size=24, color=GRAY)

multiline(s, 2, 5.8, 12, 2, [
    '"Don\'t cut more than grows back,',
    'or you destroy the resource that feeds your industry."',
], size=28, color=WHITE, spacing=8,
   bolds=[True, True])

text(s, 1, 7.5, 14, 1,
     "He called it Nachhaltigkeit. We call it sustainability.",
     size=24, color=GREEN)


# ============================================================
# SLIDE 7: That's not environmentalism
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 2.5, 14, 2,
     "THAT'S NOT",
     size=54, color=WHITE, bold=True)

text(s, 1, 4.2, 14, 1.5,
     "ENVIRONMENTALISM",
     size=54, color=WHITE, bold=True)

text(s, 1, 6.5, 14, 1,
     "That's operational resource management.",
     size=28, color=GREEN)


# ============================================================
# SLIDE 8: Industrialisation broke it
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 1, 14, 2,
     "THEN INDUSTRIALISATION",
     size=54, color=WHITE, bold=True)

text(s, 1, 2.7, 14, 1.5,
     "BROKE IT",
     size=54, color=RED, bold=True)

multiline(s, 2, 5, 12, 3, [
    "Fossil fuels removed the natural limits.",
    "You could extract faster than systems could replenish.",
    "",
    "One factory pollutes a river.",
    "A thousand factories pollute a continent.",
], size=24, color=GRAY, spacing=10)

text(s, 1, 8, 14, 0.8,
     "For 200 years, this was just called 'progress.'",
     size=22, color=DARK_GRAY, italic=True)


# ============================================================
# SLIDE 9: Costs got externalised
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 2, 14, 2,
     "COSTS GOT EXTERNALISED",
     size=54, color=WHITE, bold=True)

text(s, 1, 4.5, 14, 1,
     "The company profits.",
     size=28, color=WHITE)

text(s, 1, 5.5, 14, 1,
     "The community downstream gets sick.",
     size=28, color=RED)


# ============================================================
# SLIDE 10: 1960s-70s Rediscovery
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 1, 14, 2,
     "1960s \u2013 70s",
     size=54, color=WHITE, bold=True)

text(s, 1, 2.7, 14, 1.5,
     "THE REDISCOVERY",
     size=54, color=WHITE, bold=True)

multiline(s, 2, 5, 12, 2.5, [
    "Rachel Carson. Silent Spring. Rivers on fire.",
    "People could finally see the damage.",
    "",
    "The response: regulation.",
    "Clean Air Act. Clean Water Act. EPA.",
], size=24, color=GRAY, spacing=10)

text(s, 1, 8, 14, 0.8,
     "Straightforward: stop poisoning things. It worked for the visible stuff.",
     size=22, color=DARK_GRAY, italic=True)


# ============================================================
# SLIDE 11: Then it became a product
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 1, 14, 2,
     "1980s \u2013 2000s",
     size=54, color=WHITE, bold=True)

text(s, 1, 2.7, 14, 1.5,
     "SUSTAINABILITY BECOMES",
     size=48, color=WHITE, bold=True)

text(s, 1, 4.2, 14, 1.5,
     "A PRODUCT",
     size=48, color=AMBER, bold=True)

multiline(s, 2, 6, 12, 2.5, [
    "Companies self-regulate before regulation hits them.",
    "CSR emerges. Brundtland Report. Voluntary reporting.",
    "GRI. UN Global Compact. CDP.",
    "Reporting becomes standardized. Which means it becomes a product.",
], size=22, color=GRAY, spacing=10)


# ============================================================
# SLIDE 12: The gap
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 1.5, 14, 2,
     "THE GAP OPENS",
     size=54, color=RED, bold=True)

text(s, 1, 4, 6, 1,
     "What companies report",
     size=32, color=WHITE)

text(s, 7, 4, 2, 1,
     "vs.",
     size=32, color=DARK_GRAY)

text(s, 9, 4, 6, 1,
     "What they actually do",
     size=32, color=WHITE)

# Red divider
shape(s, MSO_SHAPE.RECTANGLE, 3, 5.5, 10, 0.05, fill=RED)

text(s, 1, 6.5, 14, 1,
     "Two different things. Increasingly so.",
     size=24, color=GRAY)


# ============================================================
# SLIDE 13: 2010s-2020s The industry
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 1, 14, 2,
     "2010s \u2013 2020s",
     size=54, color=WHITE, bold=True)

text(s, 1, 2.7, 14, 1.5,
     "THE INDUSTRY ARRIVES",
     size=54, color=WHITE, bold=True)

multiline(s, 2, 5, 12, 3, [
    "ESG becomes an investment lens. Rating agencies emerge.",
    "Supply chain questionnaires cascade downhill.",
    "CSRD. EU Taxonomy. Mandatory reporting.",
    "",
    "The sustainability industry is now worth tens of billions.",
    "Software. Consulting. Auditing. Certification.",
], size=22, color=GRAY, spacing=10)


# ============================================================
# SLIDE 14: Jim
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 2, 14, 2,
     "AND AT THE BOTTOM",
     size=48, color=WHITE, bold=True)

text(s, 1, 4.5, 14, 1,
     "Jim. 47 pages. No team. A deadline.",
     size=32, color=RED)


# ============================================================
# SLIDE 15: The lenses - title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 2, 14, 2,
     "DIFFERENT PEOPLE SEE",
     size=54, color=WHITE, bold=True)

text(s, 1, 3.7, 14, 1.5,
     "SUSTAINABILITY DIFFERENTLY",
     size=54, color=WHITE, bold=True)


# ============================================================
# SLIDE 16: Academia lens
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 1, 14, 1.5,
     "THE ACADEMIC LENS",
     size=48, color=GREEN, bold=True)

multiline(s, 2, 3.5, 12, 3, [
    "Research. Models. Peer review.",
    "",
    "Important work.",
    "But disconnected from a 15-person business.",
], size=28, color=GRAY, spacing=12)


# ============================================================
# SLIDE 17: The haters lens
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 1, 14, 1.5,
     "THE HATERS",
     size=48, color=RED, bold=True)

multiline(s, 2, 3.5, 12, 4, [
    '"ESG is woke." "Sustainability kills shareholder value."',
    "",
    "Often funded by industries that benefit from no regulation.",
    "But they have a point: the current system is performative.",
    "",
    "They just draw the wrong conclusion.",
], size=24, color=GRAY, spacing=12)


# ============================================================
# SLIDE 18: Corporate lens
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 1, 14, 1.5,
     "THE CORPORATE LENS",
     size=48, color=AMBER, bold=True)

multiline(s, 2, 3.5, 12, 3, [
    "Sustainability = reporting + compliance + brand positioning.",
    "",
    "A department, not a practice.",
], size=28, color=GRAY, spacing=12)


# ============================================================
# SLIDE 19: The alphabet soup
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 1, 14, 1.5,
     "THE ALPHABET SOUP",
     size=48, color=WHITE, bold=True)

text(s, 1, 3, 14, 1,
     "ESG  CSR  CSRD  CDP  GRI  SASB  TCFD  SFDR  VSME",
     size=28, color=DARK_GRAY)

text(s, 1, 5, 14, 1,
     "Every framework adds a layer.",
     size=28, color=GRAY)

text(s, 1, 6, 14, 1,
     "Every layer needs software, consultants, and auditors.",
     size=28, color=GRAY)


# ============================================================
# SLIDE 20: Scopes decoded
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 0.5, 14, 1,
     "IN PLAIN LANGUAGE",
     size=48, color=WHITE, bold=True)

# Scope 1
shape(s, MSO_SHAPE.RECTANGLE, 1.5, 2.3, 0.15, 1.5, fill=GREEN)
text(s, 2, 2.3, 4, 0.7, "SCOPE 1", size=28, color=GREEN, bold=True, align=PP_ALIGN.LEFT)
text(s, 2, 3, 12, 0.7, "What comes out of your building and vehicles.", size=22, color=GRAY, align=PP_ALIGN.LEFT)

# Scope 2
shape(s, MSO_SHAPE.RECTANGLE, 1.5, 4.3, 0.15, 1.5, fill=AMBER)
text(s, 2, 4.3, 4, 0.7, "SCOPE 2", size=28, color=AMBER, bold=True, align=PP_ALIGN.LEFT)
text(s, 2, 5, 12, 0.7, "Your electricity. Whatever the grid used to generate it.", size=22, color=GRAY, align=PP_ALIGN.LEFT)

# Scope 3
shape(s, MSO_SHAPE.RECTANGLE, 1.5, 6.3, 0.15, 1.5, fill=RED)
text(s, 2, 6.3, 4, 0.7, "SCOPE 3", size=28, color=RED, bold=True, align=PP_ALIGN.LEFT)
text(s, 2, 7, 12, 0.7, "Everything else. Your supply chain. Your customers' use of your product.", size=22, color=GRAY, align=PP_ALIGN.LEFT)
text(s, 2, 7.7, 12, 0.7, "This is where it gets absurd for small businesses.", size=20, color=RED, align=PP_ALIGN.LEFT)


# ============================================================
# SLIDE 21: None of these were designed for you
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 2.5, 14, 2,
     "NONE OF THESE LENSES",
     size=48, color=WHITE, bold=True)

text(s, 1, 4.2, 14, 1.5,
     "WERE DESIGNED FOR YOU",
     size=48, color=WHITE, bold=True)


# ============================================================
# SLIDE 22: Landing — 300 years
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 1.2, 14, 1.5,
     "FOR 300 YEARS,",
     size=40, color=GRAY)

text(s, 1, 2.8, 14, 1.5,
     "SUSTAINABILITY MEANT:",
     size=40, color=GRAY)

multiline(s, 2, 4.5, 12, 2, [
    '"Manage your resources',
    'so your business survives."',
], size=40, color=WHITE, spacing=8, bolds=[True, True])

# Red divider
shape(s, MSO_SHAPE.RECTANGLE, 6, 7, 4, 0.05, fill=RED)

text(s, 1, 7.5, 14, 1,
     "Somewhere along the way, it became a compliance industry.",
     size=24, color=GRAY)


# ============================================================
# SLIDE 23: Next video
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 3, 14, 2,
     "NEXT VIDEO",
     size=54, color=WHITE, bold=True)

multiline(s, 2, 5.5, 12, 2, [
    "Who does this industry actually serve?",
    "And why are the incentives designed this way?",
], size=24, color=GRAY, spacing=10)


# Save
output = r"C:\Users\User\Documents\CY\five-stacks\media\video2-the-history-v2.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Total slides: {len(prs.slides)}")
