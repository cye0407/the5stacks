"""
Build Phase 1 Video 3: THE TRAP — Why Nobody Moves First
Style: Pure black bg, white/gray text, minimal, space for images. No color coding.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA0, 0xA0, 0xA0)
DARK_GRAY = RGBColor(0x70, 0x70, 0x70)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
FONT = "Segoe UI"


def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK


def text(slide, left, top, width, height, content, size=28,
         color=WHITE, bold=False, align=PP_ALIGN.CENTER, italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = FONT
    p.alignment = align
    return txBox


def multiline(slide, left, top, width, height, lines, size=24, color=WHITE,
              align=PP_ALIGN.CENTER, spacing=14):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(spacing)
        p.alignment = align
    return txBox


# ============================================================
# SLIDE 1: Title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 2, 14, 2,
     "THE TRAP",
     size=72, color=WHITE, bold=True)

text(s, 1, 4.5, 14, 1.5,
     "Why nobody moves first.",
     size=36, color=GRAY)

text(s, 1, 8.2, 14, 0.8,
     "THE 5 STACKS  |  Video 3 of 7",
     size=16, color=DARK_GRAY)


# ============================================================
# SLIDE 2: Hook — everyone says they care
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2.5, 12, 4, [
    "Everyone says they care",
    "about the environment.",
    "",
    "Nothing actually changes.",
], size=36, color=WHITE, spacing=16)

text(s, 1, 7.5, 14, 1,
     "There's a reason for that. And it's not laziness.",
     size=24, color=GRAY)


# ============================================================
# SLIDE 3: You want to do the right thing
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2, 12, 5, [
    "You want to do the right thing.",
    "Buy the better product. Use less. Waste less.",
    "",
    "But if you do and nobody else does,",
    "you just made your life harder",
    "for nothing.",
], size=28, color=WHITE, spacing=14)


# ============================================================
# SLIDE 4: So everyone waits
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 2.5, 14, 2,
     "SO YOU WAIT",
     size=54, color=WHITE, bold=True)

text(s, 1, 5, 14, 1.5,
     "EVERYONE WAITS",
     size=54, color=WHITE, bold=True)

text(s, 1, 7.5, 14, 1,
     "Nothing happens.",
     size=28, color=GRAY)


# ============================================================
# SLIDE 5: This has a name
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2, 12, 5, [
    "This has a name.",
    "",
    "Game theory.",
    "",
    "How people make decisions when the outcome",
    "depends on what other people do.",
], size=28, color=WHITE, spacing=14)

text(s, 1, 7.5, 14, 1,
     "You already do it every day.",
     size=24, color=GRAY)


# ============================================================
# SLIDE 6: You already do this
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2.5, 12, 4, [
    "When you price a job,",
    "you think about what your competitor charges.",
    "",
    "When you invest in new equipment,",
    "you think about whether your customer will stick around.",
], size=28, color=WHITE, spacing=14)

text(s, 1, 7.5, 14, 1,
     "That's game theory.",
     size=24, color=GRAY)


# ============================================================
# SLIDE 7: Svenja's sewing company — title
# (leave space for image)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 1.5, 14, 2,
     "SVENJA RUNS",
     size=54, color=WHITE, bold=True)

text(s, 1, 3.5, 14, 1.5,
     "A SEWING COMPANY",
     size=54, color=WHITE, bold=True)

# bottom half empty for image


# ============================================================
# SLIDE 8: She could invest
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2, 12, 5, [
    "She could track her energy use.",
    "Cut her waste. Run a tighter operation.",
    "",
    "If every company in her supply chain did the same,",
    "costs would go down for everyone.",
    "The whole chain gets better.",
], size=28, color=WHITE, spacing=14)


# ============================================================
# SLIDE 9: But her competitor doesn't bother
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2, 12, 5, [
    "But her competitor doesn't bother.",
    "",
    "Spends nothing on improvement.",
    "Undercuts her on price.",
    "",
    "Svenja loses the contract.",
], size=28, color=WHITE, spacing=14)


# ============================================================
# SLIDE 10: The rational move
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 2, 14, 2,
     "THE RATIONAL MOVE:",
     size=48, color=WHITE, bold=True)

multiline(s, 2, 5, 12, 3, [
    "Don't invest. Do the bare minimum.",
    "Tick the boxes.",
    "",
    "Everyone makes this calculation.",
    "Nobody moves first.",
], size=28, color=GRAY, spacing=14)


# ============================================================
# SLIDE 11: Race to the bottom
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 3, 14, 2,
     "RACE TO THE BOTTOM",
     size=64, color=WHITE, bold=True)


# ============================================================
# SLIDE 12: The big company's response
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2, 12, 4, [
    "The big company at the top sees nobody improving.",
    "",
    "Their response:",
    "",
    '"If they won\'t do it on their own,',
    'we\'ll force them."',
], size=28, color=WHITE, spacing=14)


# ============================================================
# SLIDE 13: The questionnaire
# (leave space for image — stack of paper)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 1.5, 7, 2,
     "47 PAGES",
     size=64, color=WHITE, bold=True, align=PP_ALIGN.LEFT)

text(s, 1, 4, 7, 1,
     "Fill this out or we drop you.",
     size=28, color=GRAY, align=PP_ALIGN.LEFT)

# right half empty for image of questionnaire/paper stack


# ============================================================
# SLIDE 14: But a form doesn't fix the problem
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2, 12, 5, [
    "A questionnaire doesn't fix the problem.",
    "",
    "It doesn't change the math",
    "that made Svenja stop investing.",
    "",
    "Doing the right thing still costs more",
    "than doing nothing.",
], size=28, color=WHITE, spacing=14)


# ============================================================
# SLIDE 15: Papered over
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 2.5, 14, 2,
     "THE PROBLEM ISN'T SOLVED",
     size=48, color=WHITE, bold=True)

text(s, 1, 5.5, 14, 1.5,
     "IT'S PAPERED OVER",
     size=48, color=WHITE, bold=True)


# ============================================================
# SLIDE 16: Anchor line
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 3, 14, 2,
     "WHEN THE FIRST MOVER LOSES",
     size=48, color=WHITE, bold=True)

text(s, 1, 5.5, 14, 1.5,
     "NOBODY MOVES",
     size=54, color=WHITE, bold=True)


# ============================================================
# SLIDE 17: Not about bad people
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2.5, 12, 4, [
    "This isn't about bad people.",
    "",
    "It's about a badly designed game.",
    "",
    "The rules punish the companies",
    "that try hardest.",
], size=28, color=WHITE, spacing=14)


# ============================================================
# SLIDE 18: That's not sustainability
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 3, 14, 2,
     "THAT'S NOT SUSTAINABILITY",
     size=48, color=WHITE, bold=True)

text(s, 1, 5.5, 14, 1.5,
     "THAT'S A TRAP",
     size=48, color=WHITE, bold=True)


# ============================================================
# SLIDE 19: Next video
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 3, 14, 2,
     "NEXT",
     size=54, color=WHITE, bold=True)

multiline(s, 2, 5.5, 12, 2, [
    "If nobody moves on their own,",
    "and forms don't fix it —",
    "who actually benefits from the way things are?",
], size=24, color=GRAY, spacing=10)


# Save
output = r"C:\Users\User\Documents\CY\five-stacks\media\video3-the-trap.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Total slides: {len(prs.slides)}")
