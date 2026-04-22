"""
Build Phase 1 Video 2: The History + The Lenses — PowerPoint deck
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

INK = RGBColor(0x17, 0x17, 0x17)
TEXT = RGBColor(0x40, 0x40, 0x40)
MUTED = RGBColor(0x73, 0x73, 0x73)
SUBTLE = RGBColor(0xA3, 0xA3, 0xA3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xFA, 0xFA, 0xFA)
ACCENT = RGBColor(0x22, 0x7C, 0x5C)
LIGHT_GRAY = RGBColor(0xE5, 0xE5, 0xE5)
RED_ALERT = RGBColor(0xDC, 0x26, 0x26)
AMBER = RGBColor(0xD9, 0x77, 0x06)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
FONT = "Segoe UI"


def add_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def text(slide, left, top, width, height, content, size=28,
         color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = align
    return txBox


def multiline(slide, left, top, width, height, lines, size=24, color=TEXT,
              align=PP_ALIGN.LEFT, spacing=12, bolds=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(spacing)
        p.alignment = align
        if bolds and i < len(bolds):
            p.font.bold = bolds[i]
    return txBox


def shape(slide, stype, left, top, width, height, fill=None, line=None):
    s = slide.shapes.add_shape(stype, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(2)
    else:
        s.line.fill.background()
    return s


# ============================================================
# SLIDE 1: Title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, INK)

text(s, 1.5, 2, 13, 2,
     "THE WORD SUSTAINABILITY IS GERMAN.",
     size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

text(s, 1.5, 4, 13, 1,
     "It's 300 years old. And it had nothing to do with environmentalism.",
     size=28, color=SUBTLE, align=PP_ALIGN.CENTER)

text(s, 1.5, 6, 13, 1,
     "How did we get from resource management to a $30 billion compliance industry?",
     size=22, color=MUTED, align=PP_ALIGN.CENTER)

text(s, 1.5, 7.5, 13, 1,
     "THE 5 STACKS  |  Video 2 of 5",
     size=18, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: Indigenous land management
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 0.5, 14, 1, "Thousands of years before the word existed",
     size=40, color=INK, bold=True)

multiline(s, 1.5, 2, 6, 5, [
    "Controlled burns.",
    "Rotational farming.",
    "Seasonal harvesting.",
    "Managed forests.",
    "",
    "Every continent. Every culture.",
    "Not called 'sustainability.'",
    "Called survival.",
], size=22, color=TEXT, spacing=10)

shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 9, 2.5, 5.5, 4, fill=ACCENT)
multiline(s, 9.5, 3, 4.5, 3, [
    "You manage resources",
    "or you don't eat",
    "next year.",
    "",
    "No frameworks.",
    "No reports.",
    "Just common sense.",
], size=20, color=WHITE, align=PP_ALIGN.CENTER, spacing=8)

text(s, 1, 7.8, 14, 1,
     "Intergenerational resource management \u2014 the original sustainability.",
     size=22, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 3: Edo Japan
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 0.5, 14, 1, "Edo Japan (1603\u20131868)",
     size=44, color=INK, bold=True)

text(s, 1, 1.5, 14, 1, "250 years of near-total isolation. Nothing wasted because nothing could be replaced.",
     size=24, color=MUTED)

# Flow diagram
items = [
    (1.5, 3.5, "New clothing"),
    (5, 3.5, "Worn clothing"),
    (8.5, 3.5, "Rags"),
    (12, 3.5, "Paper"),
]

for x, y, t in items:
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 3, 1.5, fill=WHITE, line=INK)
    text(s, x + 0.2, y + 0.35, 2.6, 0.8, t, size=22, color=INK, align=PP_ALIGN.CENTER)

shape(s, MSO_SHAPE.RIGHT_ARROW, 4.6, 3.9, 0.4, 0.5, fill=INK)
shape(s, MSO_SHAPE.RIGHT_ARROW, 8.1, 3.9, 0.4, 0.5, fill=INK)
shape(s, MSO_SHAPE.RIGHT_ARROW, 11.6, 3.9, 0.4, 0.5, fill=INK)

text(s, 1, 6, 14, 1,
     "An entire economy built on 'use everything.'",
     size=28, color=INK, bold=True, align=PP_ALIGN.CENTER)

text(s, 1, 7.5, 14, 1,
     "Circularity by necessity \u2014 not by certification.",
     size=22, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4: Medieval commons
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 0.5, 14, 1, "Medieval European Commons",
     size=44, color=INK, bold=True)

text(s, 1, 1.5, 14, 1,
     "Communities set limits on extraction \u2014 not because they were green, but because overharvesting destroyed everyone's livelihood.",
     size=22, color=MUTED)

# Three pillars
pillars = [
    (2, "Shared grazing", "Limits on how many\nanimals per household"),
    (6.5, "Forestry rules", "Cut one tree,\nplant two"),
    (11, "Water rights", "Upstream users can't\npoison downstream"),
]

for x, title, desc in pillars:
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, 3.5, 3.5, 3.5, fill=WHITE, line=ACCENT)
    text(s, x + 0.2, 3.8, 3.1, 0.8, title, size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    text(s, x + 0.2, 4.8, 3.1, 1.5, desc, size=20, color=TEXT, align=PP_ALIGN.CENTER)

text(s, 1, 7.8, 14, 1,
     "Community-managed sustainability. No consultants required.",
     size=22, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5: Carlowitz 1713
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, INK)

text(s, 1, 0.8, 14, 1, "1713",
     size=80, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

text(s, 1, 2.5, 14, 1, "Hans Carl von Carlowitz",
     size=36, color=WHITE, align=PP_ALIGN.CENTER)

text(s, 1, 3.5, 14, 1, "Saxon mining administrator",
     size=24, color=SUBTLE, align=PP_ALIGN.CENTER)

shape(s, MSO_SHAPE.RECTANGLE, 6.5, 4.3, 3, 0.05, fill=MUTED)

multiline(s, 2, 5, 12, 3, [
    "The mines needed timber. They were running out.",
    "",
    '"Don\'t cut more than grows back,',
    'or you destroy the resource that feeds your industry."',
    "",
    "He called it Nachhaltigkeit.",
    "We call it sustainability.",
], size=24, color=SUBTLE, align=PP_ALIGN.CENTER, spacing=6,
   bolds=[False, False, True, True, False, True, True])

text(s, 2, 8, 12, 0.8,
     "That's not environmentalism. That's operational resource management.",
     size=22, color=ACCENT, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 6: The break — Industrialisation
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 0.5, 14, 1, "Then industrialisation broke it",
     size=44, color=INK, bold=True)

multiline(s, 1.5, 2.5, 6, 5, [
    "Fossil fuels removed the natural limits.",
    "",
    "You could extract faster than",
    "systems could replenish.",
    "",
    "One factory pollutes a river.",
    "A thousand factories pollute a continent.",
], size=22, color=TEXT, spacing=8)

shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 9, 2.5, 5.5, 3, fill=RED_ALERT)
multiline(s, 9.5, 3, 4.5, 2, [
    "Costs got externalised.",
    "The company profits.",
    "The community downstream",
    "gets sick.",
], size=20, color=WHITE, align=PP_ALIGN.CENTER, spacing=8)

text(s, 1, 7, 14, 1,
     "For 200 years, this was just called 'progress.'",
     size=28, color=MUTED, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7: 1960s-70s Rediscovery
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 0.5, 14, 1, "1960s\u201370s: The Rediscovery",
     size=44, color=INK, bold=True)

text(s, 1, 1.5, 14, 1, "People could finally see the damage.",
     size=24, color=MUTED)

multiline(s, 1.5, 3, 6.5, 4, [
    "Rachel Carson, Silent Spring (1962)",
    "Cuyahoga River on fire (1969)",
    "Visible pollution. Visible outrage.",
    "",
    "The response: regulation.",
    "Clean Air Act. Clean Water Act. EPA.",
    "",
    "Straightforward: stop poisoning things.",
    "It worked for the visible stuff.",
], size=20, color=TEXT, spacing=8)

shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 9.5, 3, 5, 3.5, fill=ACCENT)
multiline(s, 10, 3.5, 4, 2.5, [
    "This era was simple:",
    "",
    "Don't dump chemicals",
    "in the river.",
    "",
    "Everyone agreed.",
], size=20, color=WHITE, align=PP_ALIGN.CENTER, spacing=6)


# ============================================================
# SLIDE 8: 1980s-2000s Corporate capture
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 0.5, 14, 1, "1980s\u20132000s: Sustainability becomes a product",
     size=40, color=INK, bold=True)

# Timeline
decades = [
    (1.5, "1980s-90s", [
        "Companies self-regulate",
        "before regulation hits them.",
        "'CSR' emerges.",
        "Brundtland Report (1987).",
        "Sustainability becomes",
        "an industry, not a practice.",
    ]),
    (6, "2000s", [
        "GRI. UN Global Compact.",
        "CDP launches (2000).",
        "Reporting = standardized",
        "= a product to sell.",
        "Sustainability teams hired",
        "to fill out reports.",
    ]),
]

for x, decade, items in decades:
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2, 4, 5.5, fill=WHITE, line=LIGHT_GRAY)
    text(s, x + 0.3, 2.2, 3.4, 0.8, decade, size=28, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    multiline(s, x + 0.3, 3.2, 3.4, 4, items, size=18, color=TEXT, align=PP_ALIGN.CENTER, spacing=6)

shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 10.5, 2, 4.5, 5.5, fill=INK)
text(s, 10.8, 2.2, 4, 0.8, "The gap opens", size=28, color=RED_ALERT, bold=True, align=PP_ALIGN.CENTER)
multiline(s, 10.8, 3.5, 4, 3, [
    "What companies report",
    "vs.",
    "what they actually do.",
    "",
    "Two different things.",
    "Increasingly so.",
], size=20, color=WHITE, align=PP_ALIGN.CENTER, spacing=8)

text(s, 1, 8, 14, 0.8,
     "The word 'sustainability' has now left the forest and entered the boardroom.",
     size=22, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 9: 2010s-2020s The industry
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 0.5, 14, 1, "2010s\u20132020s: The industry arrives",
     size=44, color=INK, bold=True)

multiline(s, 1.5, 2, 6, 5, [
    "ESG becomes an investment lens.",
    "BlackRock. Larry Fink's letters.",
    "",
    "Rating agencies emerge:",
    "EcoVadis, Sustainalytics, MSCI ESG.",
    "",
    "Supply chain questionnaires cascade.",
    "Big companies push reporting",
    "requirements down to suppliers.",
    "",
    "Paris Agreement (2015). SDGs.",
    "CSRD, EU Taxonomy, SFDR (2020s).",
    "Mandatory reporting arrives.",
], size=20, color=TEXT, spacing=6)

shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 9, 2.5, 5.5, 4.5, fill=INK)
multiline(s, 9.5, 3, 4.5, 3.5, [
    "The sustainability industry",
    "is now worth",
    "tens of billions.",
    "",
    "Software. Consulting.",
    "Auditing. Certification.",
    "",
    "And at the bottom:",
], size=20, color=SUBTLE, align=PP_ALIGN.CENTER, spacing=6)

text(s, 9, 7.5, 5.5, 1,
     "Jim. 47 pages. No team. A deadline.",
     size=22, color=RED_ALERT, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 10: The Lenses — Overview
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 0.5, 14, 1, "Different people see sustainability differently",
     size=40, color=INK, bold=True)

# Four lens boxes
lenses = [
    (0.5, "Academia", "Research. Models. Peer review.\nImportant but disconnected\nfrom a 15-person business.", ACCENT),
    (4, "The Haters", '"ESG is woke."\nOften industry-funded.\nBut right that the system\nis performative.\nWrong conclusion.', RED_ALERT),
    (7.5, "Corporate", "Sustainability = reporting\n+ compliance + brand.\nA department, not a practice.", AMBER),
    (11, "The Alphabet Soup", "ESG, CSR, CSRD, CDP,\nGRI, SASB, TCFD, SFDR, VSME.\nEvery layer needs software\nand consultants.", INK),
]

for x, title, desc, color in lenses:
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.5, 3.8, 5, fill=WHITE, line=color)
    text(s, x + 0.2, 2.8, 3.4, 0.8, title, size=24, color=color, bold=True, align=PP_ALIGN.CENTER)
    text(s, x + 0.2, 4, 3.4, 3, desc, size=17, color=TEXT, align=PP_ALIGN.CENTER)

text(s, 1, 8, 14, 0.8,
     "None of these lenses were designed for the person who actually has to do the work.",
     size=22, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 11: The Alphabet Soup decoded
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 0.5, 14, 1, "The alphabet soup \u2014 in plain language",
     size=40, color=INK, bold=True)

# Scope definitions
scopes = [
    ("Scope 1", "What comes out of your building and vehicles.", "Direct emissions. Your gas boiler, your delivery van.", ACCENT),
    ("Scope 2", "Your electricity.", "Indirect. Whatever the grid used to generate your power.", AMBER),
    ("Scope 3", "Everything else.", "Your supply chain. Your customers' use of your product.\nThis is where it gets absurd for small businesses.", RED_ALERT),
]

y = 2
for label, desc, detail, color in scopes:
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1, y, 14, 1.8, fill=WHITE, line=color)
    text(s, 1.5, y + 0.15, 2.5, 0.8, label, size=28, color=color, bold=True)
    text(s, 4, y + 0.15, 10, 0.7, desc, size=22, color=INK, bold=True)
    text(s, 4, y + 0.85, 10, 0.7, detail, size=18, color=MUTED)
    y += 2.1

text(s, 1, 8.2, 14, 0.8,
     "Most frameworks are just different ways of asking you the same questions.",
     size=20, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 12: Landing
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, INK)

text(s, 1.5, 1.5, 13, 1.5,
     "For 300 years, sustainability meant:",
     size=36, color=SUBTLE, align=PP_ALIGN.CENTER)

text(s, 1.5, 3, 13, 1.5,
     '"Manage your resources so\nyour business survives."',
     size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

shape(s, MSO_SHAPE.RECTANGLE, 6.5, 5, 3, 0.05, fill=MUTED)

text(s, 1.5, 5.5, 13, 1.5,
     "Somewhere along the way, it became\na compliance industry worth tens of billions.",
     size=28, color=SUBTLE, align=PP_ALIGN.CENTER)

text(s, 1.5, 7.5, 13, 1,
     "Next: Who does that industry actually serve?",
     size=24, color=ACCENT, align=PP_ALIGN.CENTER)

text(s, 1.5, 8.3, 13, 0.5,
     "Subscribe  |  THE 5 STACKS",
     size=18, color=MUTED, align=PP_ALIGN.CENTER)


# Save
output = r"C:\Users\User\Documents\CY\five-stacks\media\video2-the-history.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Total slides: {len(prs.slides)}")
