"""
Build Phase 1 Video 3: "Sustainability" Isn't Sustainable
Stripped back. White on black. Space for images. No tech bro.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA0, 0xA0, 0xA0)
DARK_GRAY = RGBColor(0x70, 0x70, 0x70)
ACCENT = RGBColor(0xDC, 0x26, 0x26)  # red, used sparingly

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
FONT = "Segoe UI"


def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK


def text(slide, left, top, width, height, content, size=28,
         color=WHITE, bold=False, align=PP_ALIGN.CENTER, italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = FONT
    p.alignment = align
    return txBox


def multiline(slide, left, top, width, height, lines, size=24, color=WHITE,
              align=PP_ALIGN.CENTER, spacing=14):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(spacing)
        p.alignment = align
    return txBox


# ============================================================
# SLIDE 1: Title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 2, 14, 2,
     '"SUSTAINABILITY"',
     size=64, color=WHITE, bold=True)

text(s, 1, 4.5, 14, 1.5,
     "ISN'T SUSTAINABLE",
     size=64, color=WHITE, bold=True)

text(s, 1, 7, 14, 1,
     "And small businesses will take the blame.",
     size=24, color=GRAY)

text(s, 1, 8.2, 14, 0.8,
     "THE 5 STACKS  |  Video 3 of 5",
     size=16, color=DARK_GRAY)


# ============================================================
# SLIDE 2: You do game theory every day
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 1.5, 14, 1.5,
     "THERE'S SOMETHING CALLED",
     size=44, color=WHITE, bold=True)

text(s, 1, 3, 14, 1.5,
     "GAME THEORY",
     size=44, color=WHITE, bold=True)

multiline(s, 2, 5.5, 12, 2.5, [
    "How people make decisions when the outcome",
    "depends on what other people do.",
    "",
    "You already do this every day.",
], size=24, color=GRAY, spacing=12)


# ============================================================
# SLIDE 3: Pricing a job
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2.5, 12, 4, [
    "When you price a job,",
    "you think about what your competitor will charge.",
    "",
    "When you invest in new equipment,",
    "you think about whether your customer will stick around.",
], size=28, color=WHITE, spacing=14)

text(s, 1, 7.5, 14, 1,
     "That's game theory.",
     size=24, color=GRAY)


# ============================================================
# SLIDE 4: The trap — title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 2.5, 14, 2,
     "DOING THE RIGHT THING ALONE",
     size=48, color=WHITE, bold=True)

text(s, 1, 5, 14, 1.5,
     "COSTS YOU MORE",
     size=48, color=WHITE, bold=True)


# ============================================================
# SLIDE 5: 100 suppliers
# (leave bottom half empty for image)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 1.5, 12, 4, [
    "100 suppliers. One big customer.",
    "",
    "If they all invest in better operations,",
    "the whole supply chain improves.",
    "Costs go down. Everyone wins.",
], size=24, color=GRAY, spacing=10)

# bottom half left empty for image


# ============================================================
# SLIDE 6: But if you invest alone
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2, 12, 4, [
    "But if you invest and your competitor doesn't?",
    "",
    "They spent nothing.",
    "They undercut you on price.",
    "You lose the contract.",
], size=28, color=WHITE, spacing=14)


# ============================================================
# SLIDE 7: Bare minimum
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 2.5, 14, 2,
     "SO EVERYONE DOES",
     size=48, color=WHITE, bold=True)

text(s, 1, 5, 14, 1.5,
     "THE BARE MINIMUM",
     size=48, color=WHITE, bold=True)


# ============================================================
# SLIDE 8: The fix — questionnaire
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2, 12, 4, [
    "The big company sees this and thinks:",
    "",
    '"If they won\'t do it voluntarily,',
    'we\'ll make them."',
], size=28, color=WHITE, spacing=14)

text(s, 1, 6.5, 14, 1,
     "So they send a questionnaire.",
     size=24, color=GRAY)

text(s, 1, 7.5, 14, 1,
     "Fill this out or we drop you.",
     size=24, color=GRAY)


# ============================================================
# SLIDE 9: Papered over
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2.5, 12, 3, [
    "Forcing someone to fill out a form",
    "isn't the same as making it worthwhile",
    "to actually improve.",
], size=28, color=WHITE, spacing=14)

text(s, 1, 6.5, 14, 1,
     "The problem isn't solved. It's papered over.",
     size=24, color=GRAY, italic=True)


# ============================================================
# SLIDE 10: Who takes the risk — title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 2.5, 14, 2,
     "WHO TAKES THE RISK",
     size=54, color=WHITE, bold=True)

text(s, 1, 5.5, 14, 1.5,
     "WHO GETS THE REWARD",
     size=54, color=WHITE, bold=True)


# ============================================================
# SLIDE 11: The big company's deal
# (leave space for image on right)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 1.5, 1.5, 7, 6, [
    "Your big customer has a regulation to deal with.",
    "They need to report on their supply chain.",
    "",
    "Cheapest way? Push the work to you.",
    "You do it. For free.",
    "",
    'They write in their annual report:',
    '"We assessed 2,000 suppliers."',
    "",
    "Worth millions in reputation.",
    "Cost them almost nothing.",
], size=22, color=GRAY, align=PP_ALIGN.LEFT, spacing=8)

# right half empty for image


# ============================================================
# SLIDE 12: What you get
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2, 12, 4, [
    "You get:",
    "",
    "No help. No tools. No benefit.",
    "A deadline.",
    "And the risk that if you don't do it well enough,",
    "you lose the contract.",
], size=28, color=WHITE, spacing=14)


# ============================================================
# SLIDE 13: Rating agencies
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 1.5, 12, 5.5, [
    "Your customer pays a rating agency to rate you.",
    "You pay the same agency to see your own score.",
    "",
    "Both sides pay.",
    "The agency bears no risk.",
    "It doesn't matter if the rating",
    "makes anything more sustainable.",
    "",
    "They get paid either way.",
], size=24, color=WHITE, spacing=10)


# ============================================================
# SLIDE 14: Report vs reality
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2, 12, 4, [
    "A company with a polished sustainability report",
    "and terrible actual operations",
    "",
    "scores higher than",
    "",
    "a company with great operations",
    "and no report.",
], size=28, color=WHITE, spacing=12)


# ============================================================
# SLIDE 15: Building or flipping — title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 2.5, 14, 2,
     "ARE YOU BUILDING A BUSINESS",
     size=48, color=WHITE, bold=True)

text(s, 1, 5, 14, 1.5,
     "OR FLIPPING ONE?",
     size=48, color=WHITE, bold=True)


# ============================================================
# SLIDE 16: Quarter vs decade
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 1.5, 12, 6, [
    "Quarter to quarter:",
    "Cut corners. Push costs onto someone else.",
    "Squeeze your suppliers. Looks great on paper.",
    "",
    "Decade to decade:",
    "Know your costs. Manage your waste.",
    "Diversify your supply chain. Build something resilient.",
    "",
    "Svenja is building something her kids could take over.",
    "Those are different games.",
], size=24, color=WHITE, spacing=10)


# ============================================================
# SLIDE 17: Not the opposite of profit
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 2, 14, 2,
     "SUSTAINABILITY ISN'T THE",
     size=48, color=WHITE, bold=True)

text(s, 1, 4, 14, 1.5,
     "OPPOSITE OF PROFIT",
     size=48, color=WHITE, bold=True)

text(s, 1, 6.5, 14, 1,
     "It's what happens when you play the longer game.",
     size=28, color=GRAY)


# ============================================================
# SLIDE 18: Demands without giving back
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2, 12, 4, [
    "Your customer demands sustainability data from you.",
    "",
    "Do they give anything back?",
    "Tools? Training? Shared costs?",
    "",
    "Or do they just demand and punish?",
], size=28, color=WHITE, spacing=14)


# ============================================================
# SLIDE 19: Extraction
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 2.5, 14, 2,
     "IF IT'S ALL DEMAND",
     size=48, color=WHITE, bold=True)

text(s, 1, 4.5, 14, 1.5,
     "AND NO GIVING BACK",
     size=48, color=WHITE, bold=True)

text(s, 1, 7, 14, 1,
     "That's not a partnership. That's extraction.",
     size=24, color=GRAY)


# ============================================================
# SLIDE 20: The game right now
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2, 12, 4, [
    "Right now the game rewards:",
    "",
    "Paperwork over improvement.",
    "Extraction over cooperation.",
    "Short-term over long-term.",
    "",
    "Not because people are bad.",
    "Because the game is badly designed.",
], size=28, color=WHITE, spacing=12)


# ============================================================
# SLIDE 21: You can choose
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 2.5, 14, 2,
     "ONCE YOU SEE THE GAME",
     size=48, color=WHITE, bold=True)

text(s, 1, 5, 14, 1.5,
     "YOU CAN CHOOSE HOW TO PLAY IT",
     size=48, color=WHITE, bold=True)


# ============================================================
# SLIDE 22: Next video
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 3, 14, 2,
     "NEXT VIDEO",
     size=48, color=WHITE, bold=True)

multiline(s, 2, 5.5, 12, 2, [
    "What sustainability actually looks like",
    "when you stop playing their game",
    "and start playing your own.",
], size=24, color=GRAY, spacing=10)


# Save
output = r"C:\Users\User\Documents\CY\five-stacks\media\video3-game-theory-v4.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Total slides: {len(prs.slides)}")
