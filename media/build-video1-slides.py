"""
Build Phase 1 Video 1: The Physics — PowerPoint deck (v2)
Revised: better pacing, split dense slides, added transitions, stock vs flow
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
INK = RGBColor(0x17, 0x17, 0x17)
TEXT = RGBColor(0x40, 0x40, 0x40)
MUTED = RGBColor(0x73, 0x73, 0x73)
SUBTLE = RGBColor(0xA3, 0xA3, 0xA3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xFA, 0xFA, 0xFA)
ACCENT = RGBColor(0x22, 0x7C, 0x5C)
LIGHT_GRAY = RGBColor(0xE5, 0xE5, 0xE5)
BLUE_WATER = RGBColor(0x3B, 0x82, 0xF6)
RED_ALERT = RGBColor(0xDC, 0x26, 0x26)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

FONT = "Segoe UI"


def add_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def text(slide, left, top, width, height, content, size=28,
         color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = align
    return txBox


def multiline(slide, left, top, width, height, lines, size=24, color=TEXT, align=PP_ALIGN.LEFT, spacing=12):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(spacing)
        p.alignment = align
    return txBox


def shape(slide, stype, left, top, width, height, fill=None, line=None):
    s = slide.shapes.add_shape(stype, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(2)
    else:
        s.line.fill.background()
    return s


# ============================================================
# SLIDE 1: Title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, INK)

text(s, 1.5, 2.5, 13, 2,
     "SUSTAINABILITY ISN'T IDEOLOGY. IT'S PHYSICS.",
     size=50, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

text(s, 1.5, 5, 13, 1,
     "No opinions. No politics. Just how systems work.",
     size=28, color=SUBTLE, align=PP_ALIGN.CENTER)

text(s, 1.5, 7.5, 13, 1,
     "THE 5 STACKS  |  Video 1 of 5",
     size=18, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: Closed System
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 0.5, 14, 1, "Earth is a closed system",
     size=44, color=INK, bold=True)

text(s, 1, 1.5, 14, 1,
     "Energy comes in from the sun. That's it. Matter doesn't leave.",
     size=24, color=MUTED)

# Sun + arrow
shape(s, MSO_SHAPE.RIGHT_ARROW, 1.5, 3.5, 2.5, 0.8, fill=YELLOW)
text(s, 1.7, 3.55, 2, 0.7, "Energy in", size=18, color=INK, bold=True)

# Earth
shape(s, MSO_SHAPE.OVAL, 5, 2.8, 4.5, 4.5, fill=ACCENT, line=INK)
multiline(s, 5.5, 4.2, 3.5, 1.8,
          ["Matter stays.", "Everything we extract,", "we rearrange."],
          size=20, color=WHITE, align=PP_ALIGN.CENTER)

# Right side
multiline(s, 10.5, 3.5, 4.5, 2,
          ["No exit.", "No reset button.", "Outputs accumulate."],
          size=22, color=TEXT)


# ============================================================
# SLIDE 3: Stock vs Flow (NEW)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, INK)

text(s, 1, 0.8, 14, 1, "Stock vs. Flow",
     size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Flow box
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 2, 2.5, 5.5, 2.5, line=SUBTLE)
text(s, 2.5, 2.7, 4.5, 0.8, "FLOW", size=32, color=SUBTLE, bold=True, align=PP_ALIGN.CENTER)
text(s, 2.5, 3.5, 4.5, 1.2, "Emissions per year",
     size=24, color=WHITE, align=PP_ALIGN.CENTER)

# Stock box
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 8.5, 2.5, 5.5, 2.5, line=RED_ALERT)
text(s, 9, 2.7, 4.5, 0.8, "STOCK", size=32, color=RED_ALERT, bold=True, align=PP_ALIGN.CENTER)
text(s, 9, 3.5, 4.5, 1.2, "Total CO\u2082 in atmosphere",
     size=24, color=WHITE, align=PP_ALIGN.CENTER)

# Arrow
shape(s, MSO_SHAPE.RIGHT_ARROW, 7.7, 3.2, 0.7, 0.5, fill=SUBTLE)

# Key insight
text(s, 2, 6, 12, 1,
     "The problem isn't emissions alone. It's accumulation.",
     size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

text(s, 2, 7.5, 12, 1,
     "Even if we slow the flow, the stock keeps rising until the flow drops below what the system absorbs.",
     size=20, color=SUBTLE, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4: The Bathtub
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 0.5, 14, 1, "The Bathtub",
     size=44, color=INK, bold=True)

text(s, 1, 1.5, 14, 1,
     "Stock vs. flow \u2014 in your bathroom.",
     size=24, color=MUTED)

# Bathtub
shape(s, MSO_SHAPE.RECTANGLE, 4.5, 2.5, 6, 4.5, line=INK)
# Water (high level)
shape(s, MSO_SHAPE.RECTANGLE, 4.55, 4.0, 5.9, 2.95, fill=BLUE_WATER)

# Tap (big, red = fast)
shape(s, MSO_SHAPE.DOWN_ARROW, 6.5, 1, 1.5, 1.7, fill=RED_ALERT)
text(s, 4.2, 1.1, 2.2, 0.8, "CO\u2082 IN", size=22, color=RED_ALERT, bold=True)
text(s, 8.2, 1.1, 3.5, 0.8, "Fossil fuels, industry,\nland use", size=16, color=MUTED)

# Drain (small, green = slow)
shape(s, MSO_SHAPE.DOWN_ARROW, 7.1, 7, 0.6, 0.8, fill=ACCENT)
text(s, 5, 7.5, 2, 0.8, "CO\u2082 OUT", size=18, color=ACCENT, bold=True)
text(s, 8, 7.5, 3.5, 0.8, "Oceans, forests,\nnatural absorption", size=16, color=MUTED)

# Key message
text(s, 0.5, 3.5, 4, 1.5,
     "The tap is running\nfaster than the\ndrain can empty.",
     size=26, color=INK, bold=True)

text(s, 11, 3.5, 4.5, 1.5,
     "That's not ideology.\nThat's a bathtub\noverflowing.",
     size=26, color=TEXT)


# ============================================================
# SLIDE 5: What is a Feedback Loop?
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 0.5, 14, 1, "What is a feedback loop?",
     size=44, color=INK, bold=True)

text(s, 1, 1.8, 14, 1,
     "When the output of a system becomes an input that changes the system itself.",
     size=26, color=TEXT, align=PP_ALIGN.CENTER)

# Simple loop diagram
shape(s, MSO_SHAPE.OVAL, 5.5, 3.5, 5, 3.5, line=INK)
text(s, 6.2, 4.5, 3.5, 1.5,
     "A changes B\nB changes A\nRepeat",
     size=22, color=INK, align=PP_ALIGN.CENTER)

text(s, 1, 7.8, 14, 1,
     "Two kinds: some stabilize. Some amplify.",
     size=24, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 6: Stabilizing Feedback
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 0.5, 14, 1, "Stabilizing Feedback",
     size=44, color=ACCENT, bold=True)

text(s, 1, 1.5, 14, 1, "The system pushes back. It self-corrects.",
     size=24, color=MUTED)

# Loop boxes
steps = [
    (2, 3.5, "More heat"),
    (6.5, 3.5, "More evaporation"),
    (11, 3.5, "More clouds"),
    (6.5, 6, "Reflects sunlight\n= cools down"),
]

for x, y, t in steps:
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 3, 1.5, fill=WHITE, line=ACCENT)
    text(s, x + 0.2, y + 0.3, 2.6, 1, t, size=20, color=INK, align=PP_ALIGN.CENTER)

shape(s, MSO_SHAPE.RIGHT_ARROW, 5.2, 3.9, 1.1, 0.5, fill=ACCENT)
shape(s, MSO_SHAPE.RIGHT_ARROW, 9.7, 3.9, 1.1, 0.5, fill=ACCENT)
shape(s, MSO_SHAPE.DOWN_ARROW, 12, 5.2, 0.5, 0.8, fill=ACCENT)
shape(s, MSO_SHAPE.LEFT_ARROW, 5.2, 6.4, 1.1, 0.5, fill=ACCENT)

text(s, 1, 8, 14, 0.8, "The system has a thermostat. It tries to stay balanced.",
     size=22, color=ACCENT, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7: Amplifying Feedback
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 0.5, 14, 1, "Amplifying Feedback",
     size=44, color=RED_ALERT, bold=True)

text(s, 1, 1.5, 14, 1, "The system accelerates itself. No brakes.",
     size=24, color=MUTED)

steps = [
    (2, 3.5, "Less ice"),
    (6.5, 3.5, "Less sunlight\nreflected"),
    (11, 3.5, "More heat\nabsorbed"),
    (6.5, 6, "More ice melts"),
]

for x, y, t in steps:
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 3, 1.5, fill=WHITE, line=RED_ALERT)
    text(s, x + 0.2, y + 0.3, 2.6, 1, t, size=20, color=INK, align=PP_ALIGN.CENTER)

shape(s, MSO_SHAPE.RIGHT_ARROW, 5.2, 3.9, 1.1, 0.5, fill=RED_ALERT)
shape(s, MSO_SHAPE.RIGHT_ARROW, 9.7, 3.9, 1.1, 0.5, fill=RED_ALERT)
shape(s, MSO_SHAPE.DOWN_ARROW, 12, 5.2, 0.5, 0.8, fill=RED_ALERT)
shape(s, MSO_SHAPE.LEFT_ARROW, 5.2, 6.4, 1.1, 0.5, fill=RED_ALERT)

text(s, 1, 8, 14, 0.8,
     "This isn't a prediction. This is measured, observed, documented.",
     size=22, color=RED_ALERT, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 8: Tipping Points
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 0.5, 14, 1, "Tipping Points",
     size=44, color=INK, bold=True)

text(s, 1, 1.5, 14, 1,
     "Systems don't degrade gradually. They hold \u2014 then shift.",
     size=24, color=MUTED)

# Three bridges
# Bridge 1: holds
shape(s, MSO_SHAPE.RECTANGLE, 1.5, 4, 3.5, 0.4, fill=INK)
shape(s, MSO_SHAPE.RECTANGLE, 1.5, 4.4, 0.3, 1.5, fill=INK)
shape(s, MSO_SHAPE.RECTANGLE, 4.7, 4.4, 0.3, 1.5, fill=INK)
text(s, 1.5, 3, 3.5, 0.8, "HOLDS", size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Bridge 2: holds
shape(s, MSO_SHAPE.RECTANGLE, 6.2, 4, 3.5, 0.4, fill=INK)
shape(s, MSO_SHAPE.RECTANGLE, 6.2, 4.4, 0.3, 1.5, fill=INK)
shape(s, MSO_SHAPE.RECTANGLE, 9.4, 4.4, 0.3, 1.5, fill=INK)
text(s, 6.2, 3, 3.5, 0.8, "HOLDS", size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Bridge 3: breaks
shape(s, MSO_SHAPE.RECTANGLE, 10.9, 4.4, 0.3, 1.5, fill=INK)
shape(s, MSO_SHAPE.RECTANGLE, 14.1, 4.4, 0.3, 1.5, fill=INK)
shape(s, MSO_SHAPE.RECTANGLE, 10.9, 4, 1.5, 0.3, fill=RED_ALERT)
shape(s, MSO_SHAPE.RECTANGLE, 12.8, 4.3, 1.6, 0.3, fill=RED_ALERT)
text(s, 10.9, 3, 3.5, 0.8, "BREAKS", size=24, color=RED_ALERT, bold=True, align=PP_ALIGN.CENTER)

multiline(s, 1, 6.5, 14, 2,
          ["A rainforest can handle some clearing.",
           "Past a threshold, it becomes savanna.",
           "That transition isn't reversible on human timescales."],
          size=22, color=TEXT, align=PP_ALIGN.CENTER, spacing=8)


# ============================================================
# SLIDE 9: Bridge — "But not equally" (NEW)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, INK)

text(s, 1.5, 2.5, 13, 2,
     "These physical systems are driven\nby human activity.",
     size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

text(s, 1.5, 5.5, 13, 1,
     "But not equally.",
     size=48, color=SUBTLE, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 10: 71% — Scale
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, INK)

text(s, 1, 1.2, 14, 2, "71%",
     size=120, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

text(s, 1, 3.8, 14, 1,
     "of global industrial greenhouse gas emissions",
     size=32, color=SUBTLE, align=PP_ALIGN.CENTER)

text(s, 1, 5.2, 14, 1, "100 companies.",
     size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

multiline(s, 2, 7, 12, 1.5,
          ["The chemicals in your water. The microplastics in your food.",
           "Industrial-scale decisions by a small number of very large companies.",
           "Not Jim's tire shop. Not Svenja's sewing company."],
          size=20, color=MUTED, align=PP_ALIGN.CENTER, spacing=6)


# ============================================================
# SLIDE 11: The Real Question
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1.5, 1, 13, 1, "The physics is real.",
     size=48, color=INK, bold=True, align=PP_ALIGN.CENTER)

text(s, 1.5, 2.5, 13, 1,
     "The environment matters. That's not the debate.",
     size=32, color=TEXT, align=PP_ALIGN.CENTER)

shape(s, MSO_SHAPE.RECTANGLE, 6.5, 4, 3, 0.05, fill=LIGHT_GRAY)

text(s, 1.5, 4.5, 13, 1, "The debate is:",
     size=28, color=MUTED, align=PP_ALIGN.CENTER)

questions = [
    "Who's responsible?",
    "Who's paying?",
    "Does the system built to 'fix' it actually work?"
]

y = 5.5
for q in questions:
    text(s, 2, y, 12, 0.8, q,
         size=36, color=INK, bold=True, align=PP_ALIGN.CENTER)
    y += 1.0


# ============================================================
# SLIDE 12: Next
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, INK)

text(s, 1.5, 2, 13, 2, "Next: How did we get here?",
     size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

text(s, 1.5, 4.5, 13, 1.5,
     "From 'don't dump chemicals in the river'\nto a $30 billion compliance industry.",
     size=26, color=SUBTLE, align=PP_ALIGN.CENTER)

text(s, 1.5, 7, 13, 1,
     "Subscribe  |  THE 5 STACKS",
     size=22, color=MUTED, align=PP_ALIGN.CENTER)


# Save
output = r"C:\Users\User\Documents\CY\five-stacks\media\video1-the-physics-v2.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Total slides: {len(prs.slides)}")
