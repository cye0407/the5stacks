"""
Build Phase 1 Video 4: THE SETUP — Who Profits From the Paperwork
Style: Pure black bg, white/gray text, minimal, space for images. No color coding.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA0, 0xA0, 0xA0)
DARK_GRAY = RGBColor(0x70, 0x70, 0x70)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
FONT = "Segoe UI"


def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK


def text(slide, left, top, width, height, content, size=28,
         color=WHITE, bold=False, align=PP_ALIGN.CENTER, italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = FONT
    p.alignment = align
    return txBox


def multiline(slide, left, top, width, height, lines, size=24, color=WHITE,
              align=PP_ALIGN.CENTER, spacing=14):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(spacing)
        p.alignment = align
    return txBox


# ============================================================
# SLIDE 1: Title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 2, 14, 2,
     "THE SETUP",
     size=72, color=WHITE, bold=True)

text(s, 1, 4.5, 14, 1.5,
     "Who profits from the paperwork.",
     size=36, color=GRAY)

text(s, 1, 8.2, 14, 0.8,
     "THE 5 STACKS  |  Video 4 of 7",
     size=16, color=DARK_GRAY)


# ============================================================
# SLIDE 2: Hook — tens of billions
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2, 12, 4, [
    "There's an entire industry",
    "built on sustainability.",
    "",
    "It's worth tens of billions.",
], size=36, color=WHITE, spacing=16)

text(s, 1, 7, 14, 1,
     "Almost none of that money goes to",
     size=24, color=GRAY)

text(s, 1, 7.8, 14, 1,
     "actually making anything sustainable.",
     size=24, color=GRAY)


# ============================================================
# SLIDE 3: The everyday version — green label
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2, 12, 5, [
    "You buy the 'eco-friendly' product.",
    "It costs more.",
    "",
    "Where does the extra money go?",
], size=28, color=WHITE, spacing=14)


# ============================================================
# SLIDE 4: Where the money goes
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2, 12, 5, [
    "To the company that stuck a green label on it.",
    "To the organization that certified it.",
    "",
    "Not to the farmer.",
    "Not to the factory worker.",
    "Not to the environment.",
], size=28, color=WHITE, spacing=14)


# ============================================================
# SLIDE 5: Good intentions
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 2.5, 14, 2,
     "SOMEONE IS MAKING MONEY",
     size=48, color=WHITE, bold=True)

text(s, 1, 5, 14, 1.5,
     "OFF YOUR GOOD INTENTIONS",
     size=48, color=WHITE, bold=True)

text(s, 1, 7.5, 14, 1,
     "And it's not the people doing the work.",
     size=24, color=GRAY)


# ============================================================
# SLIDE 6: The business version — regulation
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2, 12, 5, [
    "Your big customer has a new regulation.",
    "They need to report on their supply chain.",
    "",
    "Cheapest way to do that?",
    "",
    "Push the work to you.",
], size=28, color=WHITE, spacing=14)


# ============================================================
# SLIDE 7: You do it for free
# (leave space for image on right)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 1.5, 2, 7, 5, [
    "Send you a questionnaire.",
    "You do it. For free.",
    "",
    "They write in their annual report:",
    '"We assessed 2,000 suppliers."',
    "",
    "Worth millions in reputation.",
    "Cost them almost nothing.",
], size=24, color=GRAY, align=PP_ALIGN.LEFT, spacing=10)

# right half empty for image


# ============================================================
# SLIDE 8: What you get
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 1.5, 14, 1.5,
     "WHAT YOU GET:",
     size=48, color=WHITE, bold=True)

multiline(s, 2, 4, 12, 4, [
    "No help.",
    "No tools.",
    "No benefit.",
    "A deadline.",
    "And the risk of losing the contract.",
], size=28, color=GRAY, spacing=14)


# ============================================================
# SLIDE 9: The rating agencies
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 1.5, 14, 1.5,
     "THE RATING AGENCIES",
     size=48, color=WHITE, bold=True)

multiline(s, 2, 4, 12, 4, [
    "Your customer pays them to rate you.",
    "You pay them to see your own score.",
    "",
    "Both sides pay.",
    "The agency takes no risk.",
], size=28, color=WHITE, spacing=14)


# ============================================================
# SLIDE 10: They get paid either way
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2.5, 12, 4, [
    "It doesn't matter if the rating",
    "makes anything more sustainable.",
    "",
    "They get paid either way.",
], size=28, color=WHITE, spacing=14)


# ============================================================
# SLIDE 11: Report vs reality
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 1.5, 12, 6, [
    "A company with a polished report",
    "and terrible operations",
    "",
    "scores higher than",
    "",
    "a company with great operations",
    "and no report.",
], size=28, color=WHITE, spacing=14)


# ============================================================
# SLIDE 12: Who suffers
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2.5, 12, 4, [
    "The people who designed this system",
    "don't suffer when it fails.",
    "",
    "The people at the bottom do.",
], size=28, color=WHITE, spacing=14)


# ============================================================
# SLIDE 13: Anchor line
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 2.5, 14, 2,
     "THE SYSTEM REWARDS",
     size=54, color=WHITE, bold=True)

text(s, 1, 5, 14, 1.5,
     "GOOD REPORTS",
     size=54, color=WHITE, bold=True)

text(s, 1, 7, 14, 1.5,
     "NOT GOOD OPERATIONS",
     size=54, color=WHITE, bold=True)


# ============================================================
# SLIDE 14: Not a conspiracy
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

multiline(s, 2, 2.5, 12, 4, [
    "This isn't a conspiracy.",
    "",
    "It's just incentives.",
    "",
    "The people selling the tools, the ratings,",
    "the certifications \u2014 they get paid",
    "whether anything improves or not.",
], size=28, color=WHITE, spacing=14)


# ============================================================
# SLIDE 15: Next video
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

text(s, 1, 3, 14, 2,
     "NEXT",
     size=54, color=WHITE, bold=True)

multiline(s, 2, 5.5, 12, 2, [
    "When someone tells you sustainability doesn't work \u2014",
    "are they wrong?",
    "Sort of. For about 18 months.",
], size=24, color=GRAY, spacing=10)


# Save
output = r"C:\Users\User\Documents\CY\five-stacks\media\video4-the-setup.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Total slides: {len(prs.slides)}")
